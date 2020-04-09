# COMMON IMPORTS
import os, json, csv, matplotlib.pyplot as plt
from random import randint
from datetime import datetime
from collections import Counter
from colorama import Fore


# IMPORTS RELATED TO MICROSOFT OFFICE WORD DOCUMENT
import docx
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT, WD_UNDERLINE, WD_BREAK

# IMPORTS RELATED TO MICROSOFT OFFICE POWERPOINT
import pptx
from pptx import Presentation
from pptx.util import Pt, Inches


# ==================================================================================
#   PARENT CLASS TO HANDLE REPORT TYPE, REPORT NAME, IMAGE & TABLE PROCESSING DATA
# ==================================================================================
class Head:

    def __init__(self, report_path, report_name, report_ext):
        self.write_file_path = report_path+'/'+report_name.replace(' ', '_')+datetime.now().strftime('_%d_%b_%Y')+'.'+report_ext

    # ---------------------------------------
    #   ENRICHMENT OF DATA RELATED TO IMAGE
    # ---------------------------------------
    def _process_data(self, data, field_name):
        field_list = []
        for key, val in data.items():
            if field_name in val.keys():
                field_list.append(val[field_name])
        return dict(Counter(field_list))

    # -----------------------------------------
    #   ENRICHMENT OF DATA RELATED TO TABLE
    # -----------------------------------------
    def _process_table_data(self, data, fields):
        new_dict = {}
        for key, val in data.items():
            for field in fields:
                if field in val.keys() and field not in new_dict.keys():
                    new_dict[field] = [val[field]]
                elif field in val.keys() and field in new_dict.keys():
                    new_val = new_dict[field]
                    new_val.extend([val[field]])
                    new_dict[field] = new_val
        return new_dict

# ===============================================
#  CHILD CLASS TO GENERATE REPORT IN CSV FORMAT
# ===============================================
class CSV(Head):

    def __init__(self, report_path, report_name):
        Head.__init__(self, report_path, report_name, report_ext='csv')

    # ---------------------------------------------
    #   MAIN METHOD TO INITIATE REPORT GENERATION
    # ---------------------------------------------
    def create_csv(self, json_data):
        overall_list = []
        for key, val in json_data.items():
            val.update({"@timestamp": key})
            overall_list.append(val)
        if len(overall_list) != 0:
            try:
                with open(self.write_file_path, 'w') as csvfile:
                    csvwriter = csv.DictWriter(csvfile, fieldnames=list(overall_list[0].keys()))
                    csvwriter.writeheader()
                    csvwriter.writerows(overall_list)
                return True
            except Exception as excp: return excp
        else: return False

# =================================================
#   CHILD CLASS TO GENERATE REPORT IN JSON FORMAT
# =================================================
class JSON(Head):

    def __init__(self, report_path, report_name):
        Head.__init__(self, report_path, report_name, report_ext='json')

    # ---------------------------------------------
    #   MAIN METHOD TO INITIATE REPORT GENERATION
    # ---------------------------------------------
    def create_json(self, json_obj):
        if len(json_obj) != 0:
            try:
                json_obj = json.dumps(json_obj, indent=4)
                with open(self.write_file_path, 'w') as file:
                    file.write(json_obj)
                return True
            except Exception as excp: return excp
        else: return False

# ===========================================================================
#   CHILD CLASS TO GENERATE REPORT IN MICROSOFT OFFICE WORD DOCUMENT FORMAT
# ===========================================================================
class Word(Head):

    def __init__(self, report_path, report_name, template_name, json_data, folder_path, yaml_ref):
        Head.__init__(self, report_path, report_name, report_ext='docx')
        self.rd_yml = yaml_ref
        self._doc = Document()
        self.json_data = json_data
        self._template_path = template_name
        self.folder_path = folder_path
        self._report_name = report_name

    # -----------------------------------------------------------------------------------------------------
    #   METHOD TO CREATE PARAGRAPH IN THE DOCUMENT BY ACCEPTING TEXT AND ITS PROPERTIES THROUGH PARAMETERS
    # -----------------------------------------------------------------------------------------------------
    def _text(self, text, align, underline=False, font_name=None, font_size=None, font_color=None):
        para = self._doc.add_paragraph()
        if align == 'left': para.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        elif align == 'right': para.alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT
        elif align == 'center': para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        run = para.add_run()
        run.text = text
        if underline is True: run.underline = WD_UNDERLINE.SINGLE
        if font_name is not None: run.font.name = font_name
        if font_size is not None: run.font.size = Pt(font_size)
        if font_color is not None: run.font.color.rgb = docx.shared.RGBColor(*font_color)
        del para, run

    # -----------------------------------------------------------------------------------------------------------
    #   METHOD TO ATTACH IMAGE IN THE DOCUMENT BY ACCEPTING IMAGE PATH AND ITS SPECIFICATIONS THROUGH PARAMETERS
    # -----------------------------------------------------------------------------------------------------------
    def _image(self, path, align, width, height):
        para = self._doc.add_paragraph()
        if align == 'left': para.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        elif align == 'right': para.alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT
        elif align == 'center': para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        run = para.add_run()
        run.add_picture(path, width=Inches(width), height=Inches(height))
        del para, run

    # ------------------------------------------------------------------------------------------------
    #   METHOD TO INSERT TABLE IN THE DOCUMENT BY ACCEPTING ROWS, COLS AND FIELDS THROUGH PARAMETERS
    # ------------------------------------------------------------------------------------------------
    def _table(self, rows, cols, data, fields):
        table = self._doc.add_table(rows=rows, cols=cols)
        table.style = 'Colorful List Accent 6'
        _data_dict = self._process_table_data(data=data, fields=fields)
        headings = list(_data_dict.keys())
        body = list(_data_dict.values())
        if len(headings) != 0 and len(body) != 0:
            for i in range(rows):
                if i == 0:
                    row = table.rows[i]
                    for j in range(cols):
                        row.cells[j].text = headings[j]
                else:
                    print(f'\r{Fore.WHITE}[*] Inserting Data at row : {i}', end='')
                    row = table.rows[i]
                    for j in range(cols):
                        row.cells[j].text = str(body[j][i-1])

    # --------------------------------------------------------------------------------------
    #   ADD DEFAULT COVER PAGE TO THE DOCUMENT BY INSERTING RESPECTIVE HEADINGS AND LOGO
    # --------------------------------------------------------------------------------------
    def default_cover_page(self):
        self._text(text='This is Auto generated report by the PyAlert', align='center', underline=True, font_name='verdana', font_size=14, font_color=(186, 109, 105))
        self._image(path=self.folder_path+'/report_logo.png', align='center', width=3, height=3)
        self._text(text=f'Report Name : {self._report_name}', align='center', font_name='verdana', font_size=14, font_color=(186, 109, 105))
        self._text(text=f'Report Generated at {datetime.now().strftime("%d-%b-%Y")}', align='center', font_name='arial', font_size=12, font_color=(202, 110, 12))

    # ----------------------------------------------------------------------
    #   ADD COVER PAGE WHEN USER PROVIDED INFORMATION FROM TEMPLATE FILES
    # ----------------------------------------------------------------------
    def cover_page(self):
        _cover_dict = self.rd_yml.get_word(type='cover_page', file_path=self._template_path)
        if _cover_dict is not None:
            if _cover_dict['heading'] is not None: self._text(text=_cover_dict['heading'], align='center', font_name='Biome', font_size=36, font_color=(24, 121, 191))
            if _cover_dict['subheading'] is not None: self._text(text=_cover_dict['subheading'], align='center', underline=True, font_name='verdana', font_size=14, font_color=(22, 16, 8))
            if _cover_dict['paragraph'] is not None: self._text(text=_cover_dict['paragraph'], align='center', font_name='arial', font_size=12, font_color=(22, 16, 8))

    # -----------------------------------------------------------------------------------------------------
    #   CREATE TABLE AND INSERT INTO WORD DOCUMENT BASED ON FIELDS PROVIDED BY THE USER IN TEMPLATE FILES
    # -----------------------------------------------------------------------------------------------------
    def attach_table(self):
        _table_dict = self.rd_yml.get_word(type='table', file_path=self._template_path)
        if _table_dict is not None:
            if _table_dict['heading'] is not None: self._text(text=_table_dict['heading'], align='left', underline=True, font_name='verdana', font_size=14, font_color=(12, 100, 167))
            if _table_dict['paragraph'] is not None: self._text(text=_table_dict['paragraph'], align='left', font_name='verdana', font_size=12, font_color=(8, 12, 8))
            fields = self.rd_yml.get_rule_info('table', self._template_path)['fields']
            print(f'{Fore.LIGHTBLUE_EX}[+] Observed => Total Rows : {len(self.json_data)} | Total Cols : {len(fields)}')
            self._table(rows=len(self.json_data)+1, cols=len(fields), data=self.json_data, fields=fields)

    # ----------------------------------------------------------------------------
    #   CREATE AND ATTACH IMAGE IN WORD DOCUMENT BASED ON USER PROVIDED DETAILS
    # ----------------------------------------------------------------------------
    def attach_image(self):
        _image_dict = self.rd_yml.get_word(type='charts', file_path=self._template_path)
        if _image_dict is not None:
            if _image_dict['heading'] is not None: self._text(text=_image_dict['heading'], align='left', underline=True, font_name='verdana', font_size=14, font_color=(12, 100, 167))
            _data_dict = self._process_data(self.json_data, field_name=_image_dict['field'])
            chart_obj = Chart(data=_data_dict, field_name=_image_dict['field'], folder_path=self.folder_path, save_name=randint(10, 1000))
            if _image_dict['type'] == 'pie': image_path = chart_obj.pie_chart()
            elif _image_dict['type'] == 'vbar': image_path = chart_obj.vbar_chart()
            elif _image_dict['type'] == 'hbar': image_path = chart_obj.hbar_chart()
            elif _image_dict['type'] == 'line': image_path = chart_obj.line_chart()
            self._image(path=chart_obj.save_folder, align='center', width=_image_dict['width'], height=_image_dict['height'])
            return image_path
        else: return None

    # -------------------------------------------------
    #   MAIN METHOD TO INITIATE THE REPORT GENERATION
    # -------------------------------------------------
    def create_word(self):
        if self.rd_yml.get_word(type='details', file_path=self._template_path)['type'].lower() == 'word':
            self.default_cover_page()
            self.cover_page()
            self._doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
            image_path = self.attach_image()
            self.attach_table()
            self._doc.save(self.write_file_path)
            os.remove(image_path)
            return True
        else:
            print(f'{Fore.LIGHTRED_EX}[-] Mismatched Template type. Template type should be "word" but observed template type => {self.rd_yml.get_word(type="details", file_path=self._template_path)["type"]}')
            return False

# =========================================================================
#   CHILD CLASS TO GENERATE REPORT IN MICROSOFT OFFICE POWERPOINT FORMAT
# =========================================================================
class PPT(Head):

    def __init__(self, report_path, report_name, template_name, json_data, folder_path, yaml_ref):
        Head.__init__(self, report_path, report_name, report_ext='pptx')
        self.rd_yml =yaml_ref
        self._prs = Presentation(folder_path+'/ppt_default_template.pptx')
        self.json_data = json_data
        self._template_path = template_name
        self.folder_path = folder_path
        self._report_name = report_name

    # -----------------------------------------------------------------------------------------------------
    #   METHOD TO CREATE PARAGRAPH IN THE DOCUMENT BY ACCEPTING TEXT AND ITS PROPERTIES THROUGH PARAMETERS
    # -----------------------------------------------------------------------------------------------------
    def _text(self, slide, left, top, width, height, text, font_name, font_size, font_color):
        txt_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txt_box.text_frame
        para = tf.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = pptx.dml.color.RGBColor(*font_color)

    # -----------------------------------------------------------------------------------------------------------
    #   METHOD TO ATTACH IMAGE IN THE DOCUMENT BY ACCEPTING IMAGE PATH AND ITS SPECIFICATIONS THROUGH PARAMETERS
    # -----------------------------------------------------------------------------------------------------------
    def _image(self, slide, image_path, left, top, width, height):
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))

    # ------------------------------------------------------------------------------------------------
    #   METHOD TO INSERT TABLE IN THE DOCUMENT BY ACCEPTING ROWS, COLS AND FIELDS THROUGH PARAMETERS
    # ------------------------------------------------------------------------------------------------
    def _table(self, slide, rows, cols, left, top, width, height, data, fields):
        table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
        _data_dict = self._process_table_data(data=data, fields=fields)
        headings = list(_data_dict.keys())
        body = list(_data_dict.values())
        if len(headings) != 0 and len(body) != 0:
            for i in range(rows):
                for j in range(cols):
                    if i == 0: table.cell(i, j).text = headings[j]
                    else:
                        print(f'\r{Fore.WHITE}[*] Inserting Data at row : {i}', end='')
                        table.cell(i, j).text = str(body[j][i-1])

    # --------------------------------------------------
    #   ADD NEW SLIDE TO POWERPOINT BY SPECIFYING TYPE
    # --------------------------------------------------
    def _get_slide(self, type):
        if type == 'cover': return self._prs.slides[0]
        else: return self._prs.slides.add_slide(self._prs.slide_layouts[6])

    # ----------------------------------------------------------------------
    #   ADD COVER PAGE WHEN USER PROVIDED INFORMATION FROM TEMPLATE FILES
    # ----------------------------------------------------------------------
    def cover_page(self):
        _cover_details = self.rd_yml.get_ppt(type='cover_page', file_path=self._template_path)
        slide = self._get_slide(type='cover')
        self._text(slide=slide, left=4.5, top=2, width=5, height=1.5, text=self._report_name, font_name='Biome', font_size=36, font_color=(124, 166, 203))
        self._text(slide=slide, left=5, top=3, width=5, height=1.5, text=_cover_details['heading'], font_name='verdana', font_size=25, font_color=(124, 166, 203))
        self._text(slide=slide, left=5, top=3.8, width=5, height=1.5, text=_cover_details['subheading'], font_name='arial', font_size=18, font_color=(124, 166, 203))
        self._text(slide=slide, left=5, top=4.5, width=5, height=1.5, text=f'Generated Date : {datetime.now().strftime("%d-%b-%Y")}', font_name='arial', font_size=18, font_color=(124, 166, 203))

    # -----------------------------------------------------------------------------------------------------
    #   CREATE TABLE AND INSERT INTO WORD DOCUMENT BASED ON FIELDS PROVIDED BY THE USER IN TEMPLATE FILES
    # -----------------------------------------------------------------------------------------------------
    def insert_table(self):
        slide = self._get_slide(type='page')
        _table_dict = self.rd_yml.get_word(type='table', file_path=self._template_path)
        if _table_dict is not None:
            if len(self.json_data) < 5:
                if _table_dict['heading'] is not None: self._text(slide=slide, left=1.2, top=0.5, width=3, height=1, text=_table_dict['heading'], font_name='verdana', font_size=20, font_color=(236, 181, 85))
                if _table_dict['paragraph'] is not None: self._text(slide=slide, left=1.2, top=1, width=3, height=1, text=_table_dict['paragraph'], font_name='arial', font_size=14, font_color=(255, 255, 255))
                fields = self.rd_yml.get_rule_info('table', self._template_path)['fields']
                print(f'{Fore.LIGHTBLUE_EX}[+] Observed => Total Rows : {len(self.json_data)} | Total Cols : {len(fields)}')
                self._table(slide=slide, rows=len(self.json_data)+1, cols=len(fields), data=self.json_data, fields=fields, left=1, top=1.7, width=11, height=1)
            else:
                if _table_dict['heading'] is not None: self._text(slide=slide, left=3.5, top=2, width=3, height=1, text=_table_dict['heading'], font_name='verdana', font_size=36, font_color=(236, 181, 85))
                if _table_dict['paragraph'] is not None: self._text(slide=slide, left=3.5, top=3, width=3, height=1, text=_table_dict['paragraph'], font_name='arial', font_size=25, font_color=(255, 255, 255))
                fields = self.rd_yml.get_rule_info('table', self._template_path)['fields']
                print(f'{Fore.LIGHTBLUE_EX}[+] Observed => Total Rows : {len(self.json_data)} | Total Cols : {len(fields)}')
                test_dict = {}
                for key, val in self.json_data.items():
                    if len(test_dict) < 8:
                        test_dict[key] = val
                    else:
                        test_dict[key] = val
                        slide = self._get_slide(type='page')
                        self._table(slide=slide, rows=len(test_dict) + 1, cols=len(fields), data=test_dict, fields=fields, left=1, top=0.5, width=11, height=1)
                        test_dict.clear()

    # ----------------------------------------------------------------------------
    #   CREATE AND ATTACH IMAGE IN WORD DOCUMENT BASED ON USER PROVIDED DETAILS
    # ----------------------------------------------------------------------------
    def attach_image(self):
        _image_details = self.rd_yml.get_ppt(type='charts', file_path=self._template_path)
        slide = self._get_slide(type='page')
        if _image_details['heading'] is not None: self._text(slide=slide, left=1.2, top=0.5, width=3, height=1, text=_image_details['heading'], font_name='verdana', font_size=20, font_color=(236, 181, 85))
        _data_dict = self._process_data(self.json_data, field_name=_image_details['field'])
        chart_obj = Chart(data=_data_dict, field_name=_image_details['field'], folder_path=self.folder_path, save_name=randint(10, 1000))
        if _image_details['type'] == 'pie': image_path = chart_obj.pie_chart()
        elif _image_details['type'] == 'vbar': image_path = chart_obj.vbar_chart()
        elif _image_details['type'] == 'hbar': image_path = chart_obj.hbar_chart()
        elif _image_details['type'] == 'line': image_path = chart_obj.line_chart()
        self._image(slide=slide, image_path=image_path, left=2, top=1.7, width=_image_details['width'], height=_image_details['height'])
        return image_path

    # -------------------------------------------------
    #   MAIN METHOD TO INITIATE THE REPORT GENERATION
    # -------------------------------------------------
    def create_ppt(self):
        if self.rd_yml.get_word(type='details', file_path=self._template_path)['type'].lower() == 'ppt':
            self.cover_page()
            image_path = self.attach_image()
            self.insert_table()
            self._prs.save(self.write_file_path)
            os.remove(image_path)
            return True
        else:
            print(f'{Fore.LIGHTRED_EX}[-] Mismatched Template type. Template type should be "ppt" but observed template type => {self.rd_yml.get_word(type="details", file_path=self._template_path)["type"]}')
            return False

# ==============================================
#   CLASS TO CREATE CHARTS OF MULTIPLE TYPES
# ==============================================
class Chart:

    def __init__(self, data, field_name, folder_path, save_name):
        self.data = data
        self.field_name = field_name
        self.save_folder = folder_path+'/'+str(save_name)+'.png'

    # -----------------------------
    #   METHOD TO CREATE PIE CHART
    # -----------------------------
    def pie_chart(self):
        plt.figure(figsize=(5, 5))
        plt.pie(self.data.values(), labels=self.data.keys(), startangle=150)
        plt.legend(title=self.field_name, bbox_to_anchor=(1, 0, 0.5, 1), shadow=True)
        plt.savefig(self.save_folder, dpi=400, bbox_inches='tight')
        return self.save_folder

    # ---------------------------------------
    #   METHOD TO CREATE VERTICAL BAR CHART
    # ---------------------------------------
    def vbar_chart(self):
        plt.bar(self.data.keys(), self.data.values(), color=['skyblue', 'orange', 'green', 'yellow', 'pink'], alpha=0.5, width=0.5)
        plt.xlabel(self.field_name)
        plt.ylabel('Count')
        plt.gca().spines['right'].set_visible(False)
        plt.gca().spines['top'].set_visible(False)
        for i, j in enumerate(list(self.data.values())):
            plt.text(x=i-0.4, y=j + 0.3, s=str(j))
        plt.savefig(self.save_folder, dpi=400, bbox_inches='tight')
        return self.save_folder

    # -----------------------------------------
    #   METHOD TO CREATE HORIZONTAL BAR CHART
    # -----------------------------------------
    def hbar_chart(self):
        plt.barh(list(self.data.keys()), list(self.data.values()), color=['skyblue', 'orange', 'green', 'yellow', 'pink'], alpha=0.5)
        plt.xlabel('Count')
        plt.ylabel(self.field_name)
        plt.gca().spines['right'].set_visible(False)  # remove right spine
        plt.gca().spines['top'].set_visible(False)  # remove left spine
        for i, j in enumerate(self.data.values()):
            plt.text(y=i, x=j + 0.3, s=str(j))
        plt.savefig(self.save_folder, dpi=400, bbox_inches='tight')
        return self.save_folder

    # -------------------------------
    #   METHOD TO CREATE LINE CHART
    # -------------------------------
    def line_chart(self):
        plt.plot(list(self.data.keys()), list(self.data.values()), color='c', linestyle='dashed')
        plt.xlabel(self.field_name)
        plt.ylabel('Count')
        plt.savefig(self.save_folder, dpi=400, bbox_inches='tight')
        return self.save_folder